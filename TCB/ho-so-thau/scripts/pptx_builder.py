#!/usr/bin/env python3
"""
PowerPoint Builder - Create and edit PowerPoint presentations with review tracking.

This script provides comprehensive PowerPoint creation and editing capabilities:
- Create presentations from templates or blank
- Add slides with various layouts
- Set text content (title, subtitle, bullets)
- Apply formatting standards
- Add speaker notes
- Insert tables from JSON data
- Review tracking (speaker notes, footer annotations, change log slide)

Examples:
    # Create from template with data
    python scripts/pptx_builder.py --template templates/pptx/presentation_template.pptx \\
        --data data.json --output output/presentation.pptx

    # Edit existing slide
    python scripts/pptx_builder.py --edit output/presentation.pptx --slide 3 \\
        --title "New Title" --bullets "Point 1|Point 2|Point 3"

    # Create blank presentation
    python scripts/pptx_builder.py --output output/new_presentation.pptx \\
        --title "Project Title" --subtitle "Company Name"

    # With review tracking disabled
    python scripts/pptx_builder.py --template template.pptx --data data.json \\
        --output output/presentation.pptx --no-review-mode
"""

import argparse
import json
import logging
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install with: uv pip install python-pptx")
    sys.exit(1)

try:
    import yaml
except ImportError:
    print("Error: pyyaml is required. Install with: uv pip install pyyaml")
    sys.exit(1)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


# Formatting constants
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen (33.867cm)
SLIDE_HEIGHT = Inches(7.5)    # 16:9 widescreen (19.05cm)
MARGIN = Inches(0.59)  # 1.5cm

# Font sizes
FONT_TITLE = Pt(28)
FONT_SUBTITLE = Pt(20)
FONT_BODY = Pt(18)
FONT_BULLETS = Pt(16)
FONT_FOOTER = Pt(8)

# Font family
FONT_FAMILY = 'Calibri'

# Colors
COLOR_FOOTER_GRAY = RGBColor(128, 128, 128)  # Gray for footer annotations

# Maximum bullets per slide
MAX_BULLETS = 6


class PPTXBuilder:
    """Builder class for creating and editing PowerPoint presentations."""
    
    def __init__(
        self,
        template_path: Optional[Path] = None,
        review_mode: bool = True,
        config_path: Optional[Path] = None,
        changes_path: Optional[Path] = None
    ):
        """
        Initialize the PPTX builder.
        
        Args:
            template_path: Path to template file (optional)
            review_mode: Enable review tracking (default: True)
            config_path: Path to config.yaml for author info
            changes_path: Path to changes.json for review tracking
        """
        self.review_mode = review_mode
        self.config = self._load_config(config_path)
        self.changes = self._load_changes(changes_path)
        
        # Initialize presentation
        if template_path and template_path.exists():
            self.prs = Presentation(str(template_path))
            logger.info(f"Loaded template from {template_path}")
        else:
            self.prs = Presentation()
            # Set slide dimensions to 16:9 widescreen
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT
            logger.info("Created blank presentation")
        
        # Set presentation metadata
        self._set_metadata()
    
    def _load_config(self, config_path: Optional[Path]) -> Dict:
        """Load configuration from config.yaml."""
        if config_path is None:
            config_path = Path('config.yaml')
        
        if not config_path.exists():
            logger.warning(f"Config file not found: {config_path}. Using defaults.")
            return {
                'author': {'name': 'Unknown', 'role': '', 'team': ''},
                'branding': {'company_name': 'Company'}
            }
        
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = yaml.safe_load(f)
                return config or {}
        except Exception as e:
            logger.error(f"Error loading config: {e}")
            return {
                'author': {'name': 'Unknown', 'role': '', 'team': ''},
                'branding': {'company_name': 'Company'}
            }
    
    def _load_changes(self, changes_path: Optional[Path]) -> List[Dict]:
        """Load changes manifest from changes.json."""
        if changes_path is None or not changes_path.exists():
            return []
        
        try:
            with open(changes_path, 'r', encoding='utf-8') as f:
                changes = json.load(f)
                return changes if isinstance(changes, list) else []
        except Exception as e:
            logger.error(f"Error loading changes: {e}")
            return []
    
    def _set_metadata(self):
        """Set presentation metadata from config."""
        author_name = self.config.get('author', {}).get('name', 'Unknown')
        company_name = self.config.get('branding', {}).get('company_name', 'Company')
        
        # Set core properties
        core_props = self.prs.core_properties
        core_props.author = author_name
        core_props.company = company_name
        core_props.title = self.prs.core_properties.title or 'Presentation'
    
    def add_slide(
        self,
        layout_name: str = 'content',
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
        bullets: Optional[List[str]] = None,
        notes: Optional[str] = None,
        slide_index: Optional[int] = None
    ) -> int:
        """
        Add a new slide to the presentation.
        
        Args:
            layout_name: Layout type ('title', 'content', 'two-column', 'blank', 'section')
            title: Slide title text
            subtitle: Slide subtitle text
            bullets: List of bullet points
            notes: Speaker notes text
            slide_index: Insert at specific index (None = append)
            
        Returns:
            Index of the added slide
        """
        # Map layout names to slide layouts
        layout_map = {
            'title': 0,  # Title slide layout
            'content': 1,  # Title and content layout
            'two-column': 3,  # Two content layout
            'blank': 6,  # Blank layout
            'section': 2,  # Section header layout
        }
        
        layout_idx = layout_map.get(layout_name, 1)
        
        # Get layout from presentation
        if layout_idx < len(self.prs.slide_layouts):
            layout = self.prs.slide_layouts[layout_idx]
        else:
            layout = self.prs.slide_layouts[1]  # Default to title and content
        
        # Insert slide
        if slide_index is not None:
            slide = self.prs.slides.add_slide(layout)
            # Move to desired position
            slides_list = list(self.prs.slides)
            current_idx = len(slides_list) - 1
            if slide_index < current_idx:
                # Reorder slides
                slide_id = self.prs.slides._sldIdLst[current_idx]
                self.prs.slides._sldIdLst.remove(slide_id)
                self.prs.slides._sldIdLst.insert(slide_index, slide_id)
        else:
            slide = self.prs.slides.add_slide(layout)
        
        slide_idx = len(self.prs.slides) - 1 if slide_index is None else slide_index
        
        # Set title
        if title and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._format_title(title_shape)
        
        # Set subtitle (for title slide layout)
        if subtitle:
            if layout_name == 'title' and len(slide.shapes) > 1:
                subtitle_shape = slide.shapes.placeholders[1]
                subtitle_shape.text = subtitle
                self._format_subtitle(subtitle_shape)
        
        # Set bullets
        if bullets:
            self._add_bullets(slide, bullets)
        
        # Add speaker notes
        if notes:
            self._add_speaker_notes(slide, notes)
        
        # Apply formatting
        self._format_slide(slide)
        
        logger.info(f"Added slide {slide_idx + 1} with layout '{layout_name}'")
        return slide_idx
    
    def _format_title(self, shape):
        """Format title shape with standard font."""
        if not hasattr(shape, 'text_frame'):
            return
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = FONT_FAMILY
            paragraph.font.bold = True
            paragraph.font.size = FONT_TITLE
    
    def _format_subtitle(self, shape):
        """Format subtitle shape with standard font."""
        if not hasattr(shape, 'text_frame'):
            return
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = FONT_FAMILY
            paragraph.font.size = FONT_SUBTITLE
    
    def _add_bullets(self, slide, bullets: List[str]):
        """Add bullet points to a slide."""
        # Limit to max bullets
        bullets = bullets[:MAX_BULLETS]
        
        # Find content placeholder
        content_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                phf = shape.placeholder_format
                if phf.type == 7:  # Content placeholder
                    content_shape = shape
                    break
        
        # If no content placeholder, try to find text box or create one
        if not content_shape:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    content_shape = shape
                    break
        
        if not content_shape:
            # Create text box for bullets
            left = MARGIN
            top = Inches(2)  # Below title area
            width = SLIDE_WIDTH - (2 * MARGIN)
            height = SLIDE_HEIGHT - top - MARGIN
            content_shape = slide.shapes.add_textbox(left, top, width, height)
        
        # Clear existing text
        if hasattr(content_shape, 'text_frame'):
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Add bullets
            for i, bullet_text in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet_text
                p.level = 0
                p.font.name = FONT_FAMILY
                p.font.size = FONT_BULLETS
                p.space_after = Pt(12)
    
    def _add_speaker_notes(self, slide, notes: str):
        """Add speaker notes to a slide."""
        if not slide.has_notes_slide:
            notes_slide = slide.notes_slide
        else:
            notes_slide = slide.notes_slide
        
        if notes_slide.notes_text_frame:
            # Prepend change tracking if in review mode
            if self.review_mode:
                change_info = self._get_change_info_for_slide(len(self.prs.slides) - 1)
                if change_info:
                    notes = self._format_change_tracking(change_info) + "\n\n" + notes
            
            notes_slide.notes_text_frame.text = notes
    
    def _format_slide(self, slide):
        """Apply standard formatting to a slide."""
        # Ensure slide numbers are enabled (handled by master)
        # Apply consistent margins and formatting
        
        # Add footer annotation if in review mode and slide was modified
        if self.review_mode:
            slide_idx = list(self.prs.slides).index(slide)
            change_info = self._get_change_info_for_slide(slide_idx)
            if change_info:
                self._add_footer_annotation(slide, change_info)
    
    def _add_footer_annotation(self, slide, change_info: Dict):
        """Add footer annotation to modified slides (review mode only)."""
        author = change_info.get('author', self.config.get('author', {}).get('name', 'Unknown'))
        phase = change_info.get('phase', '1')
        date = change_info.get('date', datetime.now().strftime('%Y-%m-%d'))
        
        annotation_text = f"Modified by {author} | Phase {phase} | {date}"
        
        # Create text box at bottom-left
        left = MARGIN
        top = SLIDE_HEIGHT - Inches(0.5)  # Near bottom
        width = Inches(4)
        height = Inches(0.3)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = annotation_text
        
        # Format as small gray text
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = FONT_FOOTER
        paragraph.font.color.rgb = COLOR_FOOTER_GRAY
    
    def _get_change_info_for_slide(self, slide_idx: int) -> Optional[Dict]:
        """Get change information for a specific slide from changes.json."""
        for change in self.changes:
            if change.get('slide') == slide_idx:
                return change
        return None
    
    def _format_change_tracking(self, change_info: Dict) -> str:
        """Format change tracking block for speaker notes."""
        author = change_info.get('author', self.config.get('author', {}).get('name', 'Unknown'))
        model = change_info.get('model', 'Unknown')
        phase = change_info.get('phase', '1')
        date = change_info.get('date', datetime.now().strftime('%Y-%m-%d'))
        change_type = change_info.get('change_type', 'MODIFIED')
        description = change_info.get('description', '')
        source = change_info.get('source', '')
        
        lines = [
            f"--- CHANGES (Phase {phase}) ---",
            f"{author} / {model} / {date}",
        ]
        
        if change_type == 'NEW':
            lines.append(f"- Added: {description}")
        elif change_type == 'MODIFIED':
            lines.append(f"- Modified: {description}")
        elif change_type == 'REFORMATTED':
            lines.append(f"- Reformatted: {description}")
        elif change_type == 'FIXED':
            lines.append(f"- Fixed: {description}")
        
        if source:
            lines.append(f"- Source: {source}")
        
        lines.append("---")
        
        return "\n".join(lines)
    
    def edit_slide(
        self,
        slide_index: int,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
        bullets: Optional[List[str]] = None,
        notes: Optional[str] = None
    ):
        """
        Edit an existing slide.
        
        Args:
            slide_index: Zero-based slide index
            title: New title (optional)
            subtitle: New subtitle (optional)
            bullets: New bullet points (optional, pipe-separated string or list)
            notes: New speaker notes (optional)
        """
        if slide_index < 0 or slide_index >= len(self.prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range (0-{len(self.prs.slides)-1})")
        
        slide = self.prs.slides[slide_index]
        
        # Update title
        if title is not None and slide.shapes.title:
            slide.shapes.title.text = title
            self._format_title(slide.shapes.title)
        
        # Update subtitle
        if subtitle is not None:
            # Find subtitle placeholder
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    phf = shape.placeholder_format
                    if phf.idx == 1 and len(slide.shapes) > 1:  # Usually subtitle
                        shape.text = subtitle
                        self._format_subtitle(shape)
                        break
        
        # Update bullets
        if bullets is not None:
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split('|')]
            self._add_bullets(slide, bullets)
        
        # Update speaker notes
        if notes is not None:
            self._add_speaker_notes(slide, notes)
        
        # Re-apply formatting
        self._format_slide(slide)
        
        logger.info(f"Edited slide {slide_index + 1}")
    
    def add_table(self, slide_index: int, data: List[List[str]], headers: Optional[List[str]] = None):
        """
        Add a table to a slide.
        
        Args:
            slide_index: Zero-based slide index
            data: List of rows, each row is a list of cell values
            headers: Optional header row
        """
        if slide_index < 0 or slide_index >= len(self.prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = self.prs.slides[slide_index]
        
        # Determine table dimensions
        num_rows = len(data) + (1 if headers else 0)
        num_cols = len(headers) if headers else (len(data[0]) if data else 1)
        
        # Calculate table position and size
        left = MARGIN
        top = Inches(2)  # Below title area
        width = SLIDE_WIDTH - (2 * MARGIN)
        height = SLIDE_HEIGHT - top - MARGIN - Inches(1)  # Leave space for footer
        
        # Add table
        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
        
        # Set headers
        if headers:
            for col_idx, header_text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header_text
                # Format header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.name = FONT_FAMILY
                    paragraph.font.size = FONT_BODY
        
        # Add data rows
        start_row = 1 if headers else 0
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(start_row + row_idx, col_idx)
                    cell.text = str(cell_text)
                    # Format data cells
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.name = FONT_FAMILY
                        paragraph.font.size = FONT_BODY
    
    def create_change_log_slide(self):
        """Create a 'Document Change Log' summary slide (review mode only)."""
        if not self.review_mode or not self.changes:
            return
        
        # Count changes by type
        new_count = sum(1 for c in self.changes if c.get('change_type') == 'NEW')
        modified_count = sum(1 for c in self.changes if c.get('change_type') == 'MODIFIED')
        fixed_count = sum(1 for c in self.changes if c.get('change_type') == 'FIXED')
        
        # Add summary slide
        slide_idx = self.add_slide(
            layout_name='content',
            title='Document Change Log',
            notes=f"Summary: {new_count} new slides, {modified_count} modified slides, {fixed_count} fixes"
        )
        
        slide = self.prs.slides[slide_idx]
        
        # Create table with change log data
        headers = ['Slide #', 'Change Type', 'Description', 'Author', 'Date']
        rows = []
        
        for change in self.changes:
            rows.append([
                str(change.get('slide', 'N/A') + 1),  # 1-based for display
                change.get('change_type', 'MODIFIED'),
                change.get('description', ''),
                change.get('author', self.config.get('author', {}).get('name', 'Unknown')),
                change.get('date', datetime.now().strftime('%Y-%m-%d'))
            ])
        
        # Add summary row at top
        summary_row = [
            'Summary',
            f'{new_count} new, {modified_count} modified, {fixed_count} fixed',
            '',
            '',
            ''
        ]
        rows.insert(0, summary_row)
        
        self.add_table(slide_idx, rows, headers)
        
        logger.info("Created change log slide")
    
    def build_from_data(self, data: Dict):
        """
        Build presentation from JSON data structure.
        
        Args:
            data: Dictionary with 'slides' key containing list of slide definitions
        """
        slides_data = data.get('slides', [])
        
        for slide_def in slides_data:
            slide_idx = slide_def.get('index')
            title = slide_def.get('title')
            subtitle = slide_def.get('subtitle')
            bullets = slide_def.get('bullets', [])
            notes = slide_def.get('notes', '')
            layout = slide_def.get('layout', 'content')
            
            # Add slide
            if slide_idx is not None:
                self.add_slide(
                    layout_name=layout,
                    title=title,
                    subtitle=subtitle,
                    bullets=bullets,
                    notes=notes,
                    slide_index=slide_idx
                )
            else:
                self.add_slide(
                    layout_name=layout,
                    title=title,
                    subtitle=subtitle,
                    bullets=bullets,
                    notes=notes
                )
        
        # Create change log slide if in review mode
        if self.review_mode:
            self.create_change_log_slide()
    
    def save(self, output_path: Path):
        """Save the presentation to file."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # If review mode, ensure filename ends with _REVIEW
        if self.review_mode and not output_path.stem.endswith('_REVIEW'):
            # Check if clean version exists, if so use REVIEW suffix
            clean_path = output_path.parent / f"{output_path.stem}{output_path.suffix}"
            if clean_path.exists() or '_REVIEW' not in str(output_path):
                # Create REVIEW version
                review_path = output_path.parent / f"{output_path.stem}_REVIEW{output_path.suffix}"
                self.prs.save(str(review_path))
                logger.info(f"Saved REVIEW version to {review_path}")
                return review_path
        
        self.prs.save(str(output_path))
        logger.info(f"Saved presentation to {output_path}")
        return output_path


def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        description='Create and edit PowerPoint presentations with review tracking',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    
    # Mode selection
    mode_group = parser.add_mutually_exclusive_group(required=True)
    mode_group.add_argument(
        '--template',
        type=Path,
        help='Template PowerPoint file path'
    )
    mode_group.add_argument(
        '--edit',
        type=Path,
        help='Edit existing PowerPoint file'
    )
    mode_group.add_argument(
        '--output',
        type=Path,
        help='Output file path (creates blank presentation if no template)'
    )
    
    # Data input
    parser.add_argument(
        '--data',
        type=Path,
        help='JSON data file for template filling'
    )
    
    # Slide editing options
    parser.add_argument(
        '--slide',
        type=int,
        help='Slide index (0-based) for editing'
    )
    parser.add_argument(
        '--title',
        type=str,
        help='Slide title'
    )
    parser.add_argument(
        '--subtitle',
        type=str,
        help='Slide subtitle'
    )
    parser.add_argument(
        '--bullets',
        type=str,
        help='Bullet points (pipe-separated: "Point 1|Point 2|Point 3")'
    )
    parser.add_argument(
        '--notes',
        type=str,
        help='Speaker notes text'
    )
    
    # Configuration
    parser.add_argument(
        '--config',
        type=Path,
        default=Path('config.yaml'),
        help='Path to config.yaml (default: config.yaml)'
    )
    parser.add_argument(
        '--changes',
        type=Path,
        help='Path to changes.json for review tracking'
    )
    parser.add_argument(
        '--review-mode',
        action='store_true',
        default=True,
        help='Enable review tracking (default: True)'
    )
    parser.add_argument(
        '--no-review-mode',
        dest='review_mode',
        action='store_false',
        help='Disable review tracking'
    )
    
    args = parser.parse_args()
    
    # Determine output path
    if args.output:
        output_path = args.output
    elif args.edit:
        output_path = args.edit
    elif args.template:
        # Default output name from template
        output_path = Path('output') / args.template.name
    else:
        parser.error("Must specify --output, --edit, or --template")
    
    try:
        # Initialize builder
        builder = PPTXBuilder(
            template_path=args.template,
            review_mode=args.review_mode,
            config_path=args.config,
            changes_path=args.changes
        )
        
        # Handle edit mode
        if args.edit:
            if args.slide is None:
                parser.error("--slide is required when using --edit")
            
            bullets_list = None
            if args.bullets:
                bullets_list = [b.strip() for b in args.bullets.split('|')]
            
            builder.edit_slide(
                slide_index=args.slide,
                title=args.title,
                subtitle=args.subtitle,
                bullets=bullets_list,
                notes=args.notes
            )
        
        # Handle data file mode
        elif args.data and args.data.exists():
            with open(args.data, 'r', encoding='utf-8') as f:
                data = json.load(f)
            builder.build_from_data(data)
        
        # Handle template mode (just load template, user can add slides manually)
        elif args.template:
            # Template loaded, ready for manual editing or data file
            if args.data:
                with open(args.data, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                builder.build_from_data(data)
            else:
                logger.info("Template loaded. Use --data to fill from JSON or edit programmatically.")
        
        # Save presentation
        builder.save(output_path)
        
        logger.info("Presentation created successfully!")
        
    except Exception as e:
        logger.error(f"Error: {e}", exc_info=True)
        sys.exit(1)


if __name__ == '__main__':
    main()
