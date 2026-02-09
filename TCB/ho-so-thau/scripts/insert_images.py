#!/usr/bin/env python3
"""
Insert images into Word (.docx), Excel (.xlsx), and PowerPoint (.pptx) files.

This script supports multiple insertion methods:
- Replace {{IMAGE:placeholder_name}} markers in .docx files
- Insert after specific headings in .docx files
- Insert at specific cell locations in .xlsx files
- Insert on specific slides in .pptx files

Examples:
    # Replace placeholder in Word document
    python scripts/insert_images.py --target output/file.docx --image tmp/diagrams/arch.png --placeholder architecture_diagram

    # Insert after heading in Word document
    python scripts/insert_images.py --target output/file.docx --image tmp/diagrams/arch.png --after-heading "3.2 Architecture"

    # Insert in Excel cell
    python scripts/insert_images.py --target output/file.xlsx --image tmp/diagrams/chart.png --sheet "Dashboard" --cell B2

    # Insert on PowerPoint slide
    python scripts/insert_images.py --target output/file.pptx --image tmp/diagrams/arch.png --slide 3 --centered
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

try:
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    import openpyxl
    from openpyxl.drawing.image import Image as OpenpyxlImage
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
except ImportError:
    Presentation = None


def insert_image_docx(
    target_path: Path,
    image_path: Path,
    placeholder: Optional[str] = None,
    after_heading: Optional[str] = None,
    width: float = 6.0,
) -> None:
    """
    Insert an image into a Word document.

    Args:
        target_path: Path to the .docx file
        image_path: Path to the image file
        placeholder: Placeholder name to replace (e.g., "architecture_diagram" for {{IMAGE:architecture_diagram}})
        after_heading: Heading text to insert after
        width: Image width in inches (default: 6.0)
    """
    if Document is None:
        raise ImportError("python-docx is required for .docx files. Install with: pip install python-docx")

    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    doc = Document(str(target_path))
    image_width = Inches(width)
    inserted = False

    # Method 1: Replace placeholder
    if placeholder:
        marker = f"{{{{IMAGE:{placeholder}}}}}"
        for paragraph in doc.paragraphs:
            if marker in paragraph.text:
                # Clear the paragraph and add image
                paragraph.clear()
                run = paragraph.add_run()
                run.add_picture(str(image_path), width=image_width)
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                inserted = True
                break

        # Also check in tables
        if not inserted:
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            if marker in paragraph.text:
                                paragraph.clear()
                                run = paragraph.add_run()
                                run.add_picture(str(image_path), width=image_width)
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                inserted = True
                                break
                        if inserted:
                            break
                    if inserted:
                        break
                if inserted:
                    break

    # Method 2: Insert after heading
    elif after_heading:
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip() == after_heading.strip():
                # Insert image in a new paragraph after the heading
                new_paragraph = doc.paragraphs[i].insert_paragraph_before()
                run = new_paragraph.add_run()
                run.add_picture(str(image_path), width=image_width)
                new_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                inserted = True
                break

    if not inserted:
        if placeholder:
            raise ValueError(f"Placeholder '{{{{IMAGE:{placeholder}}}}}' not found in document")
        elif after_heading:
            raise ValueError(f"Heading '{after_heading}' not found in document")

    doc.save(str(target_path))
    print(f"✓ Image inserted into {target_path.name}")


def insert_image_xlsx(
    target_path: Path,
    image_path: Path,
    sheet_name: str,
    cell: str,
    width: Optional[int] = None,
    height: Optional[int] = None,
) -> None:
    """
    Insert an image into an Excel worksheet.

    Args:
        target_path: Path to the .xlsx file
        image_path: Path to the image file
        sheet_name: Name of the worksheet
        cell: Cell reference (e.g., "B2")
        width: Image width in pixels (optional, maintains aspect ratio if not specified)
        height: Image height in pixels (optional, maintains aspect ratio if not specified)
    """
    if openpyxl is None:
        raise ImportError("openpyxl is required for .xlsx files. Install with: pip install openpyxl")

    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    workbook = openpyxl.load_workbook(str(target_path))

    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Sheet '{sheet_name}' not found. Available sheets: {', '.join(workbook.sheetnames)}")

    worksheet = workbook[sheet_name]

    # Create image object
    img = OpenpyxlImage(str(image_path))

    # Set dimensions if specified
    if width:
        img.width = width
    if height:
        img.height = height

    # Add image to cell
    worksheet.add_image(img, cell)

    workbook.save(str(target_path))
    print(f"✓ Image inserted into {target_path.name} at sheet '{sheet_name}', cell {cell}")


def insert_image_pptx(
    target_path: Path,
    image_path: Path,
    slide_index: int,
    centered: bool = False,
    left: Optional[float] = None,
    top: Optional[float] = None,
    width: Optional[float] = None,
    height: Optional[float] = None,
) -> None:
    """
    Insert an image into a PowerPoint slide.

    Args:
        target_path: Path to the .pptx file
        image_path: Path to the image file
        slide_index: Slide index (0-based)
        centered: Center the image on the slide
        left: Left position in inches (optional)
        top: Top position in inches (optional)
        width: Image width in inches (optional)
        height: Image height in inches (optional)
    """
    if Presentation is None:
        raise ImportError("python-pptx is required for .pptx files. Install with: pip install python-pptx")

    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    prs = Presentation(str(target_path))

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range. Document has {len(prs.slides)} slides (0-{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate position and size
    # Get image dimensions for aspect ratio calculation
    try:
        from PIL import Image as PILImage
        pil_img = PILImage.open(str(image_path))
        img_aspect_ratio = pil_img.height / pil_img.width
    except ImportError:
        # Fallback if PIL not available - assume 4:3 ratio
        img_aspect_ratio = 0.75
        print("Warning: PIL not available, using default aspect ratio. Install Pillow for accurate sizing.", file=sys.stderr)
    except Exception as e:
        # Fallback on any error
        img_aspect_ratio = 0.75
        print(f"Warning: Could not read image dimensions: {e}. Using default aspect ratio.", file=sys.stderr)

    if centered:
        # Center the image
        if width is None:
            # Default width: 80% of slide width (convert EMU to inches)
            slide_width_inches = slide_width / 914400
            img_width = PptxInches(slide_width_inches * 0.8)
        else:
            img_width = PptxInches(width)

        if height is None:
            # Maintain aspect ratio
            img_height = PptxInches(img_width.inches * img_aspect_ratio)
        else:
            img_height = PptxInches(height)

        # Center position (convert EMU to inches)
        slide_width_inches = slide_width / 914400
        slide_height_inches = slide_height / 914400
        left_pos = PptxInches((slide_width_inches - img_width.inches) / 2)
        top_pos = PptxInches((slide_height_inches - img_height.inches) / 2)
    else:
        # Use specified positions or defaults
        left_pos = PptxInches(left) if left is not None else PptxInches(1.0)
        top_pos = PptxInches(top) if top is not None else PptxInches(1.0)

        if width is None:
            img_width = PptxInches(6.0)  # Default width
        else:
            img_width = PptxInches(width)

        if height is None:
            # Maintain aspect ratio
            img_height = PptxInches(img_width.inches * img_aspect_ratio)
        else:
            img_height = PptxInches(height)

    # Add image to slide
    slide.shapes.add_picture(str(image_path), left_pos, top_pos, img_width, img_height)

    prs.save(str(target_path))
    print(f"✓ Image inserted into {target_path.name} on slide {slide_index + 1}")


def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Insert images into Word, Excel, or PowerPoint files",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Replace placeholder in Word document
  %(prog)s --target output/file.docx --image tmp/diagrams/arch.png --placeholder architecture_diagram

  # Insert after heading in Word document
  %(prog)s --target output/file.docx --image tmp/diagrams/arch.png --after-heading "3.2 Architecture" --width 5.0

  # Insert in Excel cell
  %(prog)s --target output/file.xlsx --image tmp/diagrams/chart.png --sheet "Dashboard" --cell B2 --excel-width 400 --excel-height 300

  # Insert on PowerPoint slide (centered)
  %(prog)s --target output/file.pptx --image tmp/diagrams/arch.png --slide 3 --centered

  # Insert on PowerPoint slide (custom position)
  %(prog)s --target output/file.pptx --image tmp/diagrams/arch.png --slide 3 --left 1.0 --top 2.0 --pptx-width 8.0 --pptx-height 6.0
        """
    )

    parser.add_argument(
        "--target",
        type=Path,
        required=True,
        help="Target document file (.docx, .xlsx, or .pptx)"
    )
    parser.add_argument(
        "--image",
        type=Path,
        required=True,
        help="Path to the image file to insert"
    )

    # Word-specific options
    word_group = parser.add_argument_group("Word document options (.docx)")
    word_group.add_argument(
        "--placeholder",
        type=str,
        help="Placeholder name to replace (e.g., 'architecture_diagram' for {{IMAGE:architecture_diagram}})"
    )
    word_group.add_argument(
        "--after-heading",
        type=str,
        help="Insert image after this heading text"
    )
    word_group.add_argument(
        "--width",
        type=float,
        default=6.0,
        help="Image width in inches (default: 6.0)"
    )

    # Excel-specific options
    excel_group = parser.add_argument_group("Excel document options (.xlsx)")
    excel_group.add_argument(
        "--sheet",
        type=str,
        help="Worksheet name"
    )
    excel_group.add_argument(
        "--cell",
        type=str,
        help="Cell reference (e.g., 'B2')"
    )
    excel_group.add_argument(
        "--excel-width",
        type=int,
        help="Image width in pixels (Excel only)"
    )
    excel_group.add_argument(
        "--excel-height",
        type=int,
        help="Image height in pixels (Excel only)"
    )

    # PowerPoint-specific options
    pptx_group = parser.add_argument_group("PowerPoint document options (.pptx)")
    pptx_group.add_argument(
        "--slide",
        type=int,
        help="Slide index (0-based)"
    )
    pptx_group.add_argument(
        "--centered",
        action="store_true",
        help="Center the image on the slide"
    )
    pptx_group.add_argument(
        "--left",
        type=float,
        help="Left position in inches (PowerPoint only)"
    )
    pptx_group.add_argument(
        "--top",
        type=float,
        help="Top position in inches (PowerPoint only)"
    )
    pptx_group.add_argument(
        "--pptx-width",
        type=float,
        help="Image width in inches (PowerPoint only)"
    )
    pptx_group.add_argument(
        "--pptx-height",
        type=float,
        help="Image height in inches (PowerPoint only)"
    )

    args = parser.parse_args()

    # Validate target file exists
    if not args.target.exists():
        print(f"Error: Target file not found: {args.target}", file=sys.stderr)
        sys.exit(1)

    # Auto-detect format from extension
    ext = args.target.suffix.lower()
    
    try:
        if ext == ".docx":
            if not args.placeholder and not args.after_heading:
                print("Error: For .docx files, either --placeholder or --after-heading must be specified", file=sys.stderr)
                sys.exit(1)
            insert_image_docx(
                args.target,
                args.image,
                placeholder=args.placeholder,
                after_heading=args.after_heading,
                width=args.width
            )

        elif ext == ".xlsx":
            if not args.sheet or not args.cell:
                print("Error: For .xlsx files, both --sheet and --cell must be specified", file=sys.stderr)
                sys.exit(1)
            insert_image_xlsx(
                args.target,
                args.image,
                sheet_name=args.sheet,
                cell=args.cell,
                width=args.excel_width,
                height=args.excel_height
            )

        elif ext == ".pptx":
            if args.slide is None:
                print("Error: For .pptx files, --slide must be specified", file=sys.stderr)
                sys.exit(1)
            insert_image_pptx(
                args.target,
                args.image,
                slide_index=args.slide,
                centered=args.centered,
                left=args.left,
                top=args.top,
                width=args.pptx_width,
                height=args.pptx_height
            )

        else:
            print(f"Error: Unsupported file format: {ext}. Supported formats: .docx, .xlsx, .pptx", file=sys.stderr)
            sys.exit(1)

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
