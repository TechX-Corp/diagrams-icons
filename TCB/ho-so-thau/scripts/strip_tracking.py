#!/usr/bin/env python3
"""
Strip review tracking artifacts from REVIEW versions of documents.

This script removes all review tracking elements (comments, highlights, change logs,
audit columns, etc.) from REVIEW versions of documents and creates clean versions
suitable for client delivery.

Usage:
    python scripts/strip_tracking.py output/proposal_REVIEW.docx
    python scripts/strip_tracking.py output/*_REVIEW.*
    python scripts/strip_tracking.py output/proposal_REVIEW.docx --output output/proposal.docx
"""

import argparse
import glob
import logging
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.oxml import parse_xml
    from lxml import etree
except ImportError:
    Document = None
    logging.warning("python-docx/lxml not available. .docx stripping will be disabled.")

try:
    from openpyxl import load_workbook
    from openpyxl.styles import PatternFill
except ImportError:
    load_workbook = None
    logging.warning("openpyxl not available. .xlsx stripping will be disabled.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logging.warning("python-pptx not available. .pptx stripping will be disabled.")


logging.basicConfig(
    level=logging.WARNING,
    format='%(levelname)s: %(message)s'
)
logger = logging.getLogger(__name__)


def strip_docx_tracking(review_path: Path, clean_path: Path) -> Dict[str, int]:
    """
    Strip review tracking from a Word document.
    
    Args:
        review_path: Path to the REVIEW version
        clean_path: Path to save the clean version
        
    Returns:
        Dictionary with counts of removed elements
    """
    if Document is None:
        logger.error("python-docx is not installed. Cannot process .docx files.")
        return {}
    
    stats = {
        'comments': 0,
        'highlights': 0,
        'update_markers': 0,
        'change_log_sections': 0
    }
    
    try:
        doc = Document(str(review_path))
        
        # Remove comments from document.xml
        # Comments are stored in comments.xml and referenced in document.xml
        # We need to remove comment ranges and comment references
        
        # Get the document part XML
        doc_part = doc.part
        doc_xml = doc_part.element
        
        # Namespace for Word processing
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        
        # Remove comment range starts
        comment_range_starts = doc_xml.xpath('.//w:commentRangeStart', namespaces=ns)
        for elem in comment_range_starts:
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)
                stats['comments'] += 1
        
        # Remove comment range ends
        comment_range_ends = doc_xml.xpath('.//w:commentRangeEnd', namespaces=ns)
        for elem in comment_range_ends:
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)
                stats['comments'] += 1
        
        # Remove comment reference runs (these are the visible comment markers)
        comment_refs = doc_xml.xpath('.//w:commentReference', namespaces=ns)
        for elem in comment_refs:
            # Remove the run containing the comment reference
            run = elem.getparent()
            if run is not None and run.tag == qn('w:r'):
                parent = run.getparent()
                if parent is not None:
                    parent.remove(run)
                    stats['comments'] += 1
        
        # Remove highlighting from all runs
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run._element.rPr is not None:
                    highlight_elem = run._element.rPr.find(qn('w:highlight'))
                    if highlight_elem is not None:
                        run._element.rPr.remove(highlight_elem)
                        stats['highlights'] += 1
        
        # Also check tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run._element.rPr is not None:
                                highlight_elem = run._element.rPr.find(qn('w:highlight'))
                                if highlight_elem is not None:
                                    run._element.rPr.remove(highlight_elem)
                                    stats['highlights'] += 1
        
        # Remove section-level update markers (italic gray text matching pattern)
        # Pattern: "Updated by.*using.*on.*Phase"
        update_marker_pattern = re.compile(r'Updated by.*using.*on.*Phase', re.IGNORECASE)
        
        for paragraph in doc.paragraphs:
            # Check if paragraph contains update marker pattern
            para_text = paragraph.text
            if update_marker_pattern.search(para_text):
                # Check if it's italic and gray (likely an update marker)
                # Remove the entire paragraph if it matches
                if para_text.strip():
                    # Check runs for italic formatting
                    is_marker = False
                    for run in paragraph.runs:
                        if run.italic or run.font.color and 'gray' in str(run.font.color).lower():
                            is_marker = True
                            break
                    
                    if is_marker and update_marker_pattern.search(para_text):
                        # Remove this paragraph
                        p_elem = paragraph._element
                        p_elem.getparent().remove(p_elem)
                        stats['update_markers'] += 1
        
        # Remove Change Log section
        # Find heading containing "Change Log" and delete from that heading to end or next H1
        change_log_heading = None
        change_log_start_idx = None
        
        for i, paragraph in enumerate(doc.paragraphs):
            para_text = paragraph.text.strip()
            # Check if it's a heading and contains "Change Log"
            if 'Change Log' in para_text or 'change log' in para_text.lower():
                # Check if it's styled as a heading (has heading style)
                if paragraph.style.name.startswith('Heading'):
                    change_log_heading = paragraph
                    change_log_start_idx = i
                    break
        
        if change_log_start_idx is not None:
            # Remove all paragraphs from this point to the end
            # Or until we hit the next H1 heading
            paragraphs_to_remove = []
            for i in range(change_log_start_idx, len(doc.paragraphs)):
                para = doc.paragraphs[i]
                # Stop if we hit another H1 heading (but not the Change Log heading itself)
                if i > change_log_start_idx and para.style.name == 'Heading 1':
                    break
                paragraphs_to_remove.append(para._element)
            
            # Remove paragraphs in reverse order to avoid index issues
            for p_elem in reversed(paragraphs_to_remove):
                p_elem.getparent().remove(p_elem)
                stats['change_log_sections'] += 1
        
        # Remove comments.xml part and its relationships
        # Comments are stored in a separate part, we need to remove the relationship
        try:
            # Find and remove comment relationships from the relationships part
            rels_part = doc_part.rels
            rels_to_remove = []
            for rel_id, rel in rels_part.items():
                if 'comments' in rel.target_ref.lower() or rel.reltype.endswith('comments'):
                    rels_to_remove.append(rel_id)
            
            for rel_id in rels_to_remove:
                rels_part.remove(rel_id)
        except Exception as e:
            logger.debug(f"Could not remove comments relationship: {e}")
        
        # Save the clean document
        doc.save(str(clean_path))
        
    except Exception as e:
        logger.error(f"Error processing {review_path}: {e}")
        raise
    
    return stats


def strip_xlsx_tracking(review_path: Path, clean_path: Path) -> Dict[str, int]:
    """
    Strip review tracking from an Excel workbook.
    
    Args:
        review_path: Path to the REVIEW version
        clean_path: Path to save the clean version
        
    Returns:
        Dictionary with counts of removed elements
    """
    if load_workbook is None:
        logger.error("openpyxl is not installed. Cannot process .xlsx files.")
        return {}
    
    stats = {
        'audit_columns': 0,
        'row_colors': 0,
        'change_log_sheet': 0,
        'cell_comments': 0
    }
    
    try:
        wb = load_workbook(str(review_path))
        
        # Track columns to delete across all sheets
        audit_column_names = ['_Status', '_Change_Desc', '_Updated_By']
        
        # Process each sheet
        for sheet_name in list(wb.sheetnames):
            ws = wb[sheet_name]
            
            # Delete _Change_Log sheet
            if sheet_name == '_Change_Log':
                wb.remove(ws)
                stats['change_log_sheet'] += 1
                continue
            
            # Find and delete audit columns
            columns_to_delete = []
            header_row = 1
            
            # Check header row for audit column names
            for col_idx, cell in enumerate(ws[header_row], 1):
                if cell.value in audit_column_names:
                    columns_to_delete.append(col_idx)
            
            # Delete columns in reverse order to maintain indices
            for col_idx in reversed(sorted(columns_to_delete)):
                ws.delete_cols(col_idx)
                stats['audit_columns'] += 1
            
            # Remove row background colors (only green #E2EFDA and yellow #FFF2CC)
            no_fill = PatternFill()
            
            for row in ws.iter_rows():
                for cell in row:
                    if cell.fill and cell.fill.fill_type == 'solid':
                        # Get color value - can be RGB string or hex
                        fill_color = None
                        if hasattr(cell.fill.start_color, 'rgb'):
                            fill_color = cell.fill.start_color.rgb
                        elif hasattr(cell.fill.start_color, 'index'):
                            # Indexed color, skip (not our tracking colors)
                            continue
                        else:
                            fill_color = str(cell.fill.start_color)
                        
                        # Check if it's the tracking colors (case-insensitive, with or without #)
                        fill_color_str = str(fill_color).upper().replace('#', '')
                        if fill_color_str == 'E2EFDA' or fill_color_str == 'FFF2CC':
                            cell.fill = no_fill
                            stats['row_colors'] += 1
            
            # Remove all cell comments
            for row in ws.iter_rows():
                for cell in row:
                    if cell.comment:
                        cell.comment = None
                        stats['cell_comments'] += 1
        
        # Save the clean workbook
        wb.save(str(clean_path))
        
    except Exception as e:
        logger.error(f"Error processing {review_path}: {e}")
        raise
    
    return stats


def strip_pptx_tracking(review_path: Path, clean_path: Path) -> Dict[str, int]:
    """
    Strip review tracking from a PowerPoint presentation.
    
    Args:
        review_path: Path to the REVIEW version
        clean_path: Path to save the clean version
        
    Returns:
        Dictionary with counts of removed elements
    """
    if Presentation is None:
        logger.error("python-pptx is not installed. Cannot process .pptx files.")
        return {}
    
    stats = {
        'speaker_note_changes': 0,
        'footer_annotations': 0,
        'change_log_slides': 0
    }
    
    try:
        prs = Presentation(str(review_path))
        
        # Pattern for change tracking block in speaker notes
        change_block_pattern = re.compile(
            r'---\s*CHANGES.*?---',
            re.DOTALL | re.IGNORECASE
        )
        
        # Pattern for footer annotation
        footer_pattern = re.compile(r'Modified by.*Phase', re.IGNORECASE)
        
        slides_to_remove = []
        
        for slide_idx, slide in enumerate(prs.slides):
            # Remove slide titled "Document Change Log"
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if 'Change Log' in title_text or 'change log' in title_text.lower():
                    slides_to_remove.append(slide_idx)
                    stats['change_log_slides'] += 1
                    continue
            
            # Process speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text
                    
                    # Remove change tracking block
                    if change_block_pattern.search(notes_text):
                        # Remove the change block from notes
                        cleaned_text = change_block_pattern.sub('', notes_text).strip()
                        notes_slide.notes_text_frame.text = cleaned_text
                        stats['speaker_note_changes'] += 1
            
            # Remove footer annotation text boxes
            # Check all shapes for footer annotation pattern
            shapes_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if footer_pattern.search(shape.text):
                        # Check if it's positioned at the bottom (likely a footer)
                        # PowerPoint slide height is typically around 6858000 EMU (19.05cm)
                        if shape.top > 6000000:  # Near bottom of slide
                            shapes_to_remove.append(shape)
                            stats['footer_annotations'] += 1
            
            # Remove shapes (in reverse to avoid index issues)
            for shape in reversed(shapes_to_remove):
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Remove slides (in reverse order to maintain indices)
        for slide_idx in reversed(sorted(slides_to_remove)):
            try:
                slide_id = prs.slides._sldIdLst[slide_idx]
                rId = slide_id.rId
                # Remove the relationship
                prs.part.drop_rel(rId)
                # Remove from slide list
                prs.slides._sldIdLst.remove(slide_id)
            except Exception as e:
                logger.warning(f"Could not remove slide {slide_idx}: {e}")
        
        # Save the clean presentation
        prs.save(str(clean_path))
        
    except Exception as e:
        logger.error(f"Error processing {review_path}: {e}")
        raise
    
    return stats


def get_clean_path(review_path: Path, output_path: Optional[Path] = None) -> Path:
    """
    Determine the clean output path from the review path.
    
    Args:
        review_path: Path to the REVIEW version
        output_path: Optional explicit output path
        
    Returns:
        Path for the clean version
    """
    if output_path:
        return Path(output_path)
    
    # Remove _REVIEW from filename
    stem = review_path.stem
    if stem.endswith('_REVIEW'):
        clean_stem = stem[:-7]  # Remove '_REVIEW'
    else:
        clean_stem = stem
    
    return review_path.parent / f"{clean_stem}{review_path.suffix}"


def strip_tracking(review_path: Path, output_path: Optional[Path] = None) -> Dict[str, int]:
    """
    Strip review tracking from a document based on its file extension.
    
    Args:
        review_path: Path to the REVIEW version
        output_path: Optional explicit output path
        
    Returns:
        Dictionary with counts of removed elements
    """
    review_path = Path(review_path)
    
    if not review_path.exists():
        logger.error(f"File not found: {review_path}")
        return {}
    
    clean_path = get_clean_path(review_path, output_path)
    suffix = review_path.suffix.lower()
    
    if suffix == '.docx':
        return strip_docx_tracking(review_path, clean_path)
    elif suffix == '.xlsx':
        return strip_xlsx_tracking(review_path, clean_path)
    elif suffix == '.pptx':
        return strip_pptx_tracking(review_path, clean_path)
    else:
        logger.warning(f"Unsupported file type: {suffix}. Skipping {review_path}")
        return {}


def print_summary(file_path: Path, stats: Dict[str, int]):
    """
    Print a summary of what was stripped from a file.
    
    Args:
        file_path: Path to the processed file
        stats: Dictionary with counts of removed elements
    """
    if not stats:
        return
    
    print(f"\n{file_path.name}:")
    total = 0
    
    for key, count in stats.items():
        if count > 0:
            key_display = key.replace('_', ' ').title()
            print(f"  - Removed {count} {key_display}")
            total += count
    
    if total == 0:
        print("  - No tracking artifacts found")


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description='Strip review tracking artifacts from REVIEW versions of documents',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python scripts/strip_tracking.py output/proposal_REVIEW.docx
  python scripts/strip_tracking.py output/*_REVIEW.*
  python scripts/strip_tracking.py output/proposal_REVIEW.docx --output output/proposal.docx
        """
    )
    
    parser.add_argument(
        'input_paths',
        type=str,
        nargs='+',
        help='Path(s) to REVIEW version file(s) or glob pattern(s)'
    )
    
    parser.add_argument(
        '--output',
        type=str,
        default=None,
        help='Output path (only valid for single file input)'
    )
    
    parser.add_argument(
        '--verbose',
        action='store_true',
        help='Enable verbose logging'
    )
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.INFO)
    
    # Expand glob patterns
    file_paths = []
    for pattern in args.input_paths:
        expanded = glob.glob(pattern)
        if expanded:
            file_paths.extend(expanded)
        else:
            # If glob didn't match, treat as literal path
            file_paths.append(pattern)
    
    if not file_paths:
        logger.error("No files found matching the provided patterns")
        sys.exit(1)
    
    # Validate --output usage
    if args.output and len(file_paths) > 1:
        logger.error("--output can only be used with a single input file")
        sys.exit(1)
    
    # Process each file
    total_stats = {}
    processed_count = 0
    
    for file_path_str in file_paths:
        file_path = Path(file_path_str)
        
        if not file_path.exists():
            logger.warning(f"File not found: {file_path}. Skipping.")
            continue
        
        try:
            output_path = Path(args.output) if args.output else None
            stats = strip_tracking(file_path, output_path)
            
            if stats:
                print_summary(file_path, stats)
                processed_count += 1
                
                # Aggregate stats
                for key, count in stats.items():
                    total_stats[key] = total_stats.get(key, 0) + count
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            continue
    
    # Print overall summary
    if processed_count > 0:
        print(f"\n{'='*60}")
        print(f"Processed {processed_count} file(s)")
        if total_stats:
            print("Total removed:")
            for key, count in total_stats.items():
                if count > 0:
                    key_display = key.replace('_', ' ').title()
                    print(f"  - {count} {key_display}")
        print(f"{'='*60}")
    else:
        logger.warning("No files were processed successfully")


if __name__ == "__main__":
    main()
