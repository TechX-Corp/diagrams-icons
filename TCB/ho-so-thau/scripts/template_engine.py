#!/usr/bin/env python3
"""
Template engine for filling Word (.docx), Excel (.xlsx), and PowerPoint (.pptx) templates.

This script replaces {{placeholder}} markers in templates with data from JSON files.
It supports:
- Simple value replacement: {{company_name}} -> "TechX Solutions"
- Nested placeholders in table cells
- List data with row replication (e.g., pricing line items)
- Preserves formatting of placeholder text
- Leaves {{IMAGE:name}} markers untouched (handled by insert_images.py)

Examples:
    # Direct mode: specify template and data files
    python scripts/template_engine.py --template templates/word/proposal_template.docx --data data.json --output output/proposal.docx

    # Config mode: use template_config.yaml
    python scripts/template_engine.py --config templates/template_config.yaml --type proposal --data data.json
"""

import argparse
import json
import re
import sys
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from docx import Document as DocumentType
    import openpyxl.workbook
    from pptx import Presentation as PresentationType

try:
    import yaml
except ImportError:
    yaml = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from jinja2 import Template as JinjaTemplate
except ImportError:
    JinjaTemplate = None


# Regex pattern to match {{placeholder}} markers (but not {{IMAGE:...}})
PLACEHOLDER_PATTERN = re.compile(r'\{\{([^}]+)\}\}')
IMAGE_PLACEHOLDER_PATTERN = re.compile(r'\{\{IMAGE:([^}]+)\}\}')


def load_data(data_path: Path) -> Dict[str, Any]:
    """
    Load data from JSON file.

    Args:
        data_path: Path to JSON data file

    Returns:
        Dictionary containing the data
    """
    if not data_path.exists():
        raise FileNotFoundError(f"Data file not found: {data_path}")

    with open(data_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def load_template_config(config_path: Path) -> Dict[str, Any]:
    """
    Load template configuration from YAML file.

    Args:
        config_path: Path to template_config.yaml

    Returns:
        Dictionary containing template configuration
    """
    if yaml is None:
        raise ImportError("PyYAML is required for config mode. Install with: pip install pyyaml")

    if not config_path.exists():
        warnings.warn(f"Template config not found: {config_path}. Skipping validation.")
        return {}

    with open(config_path, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f) or {}


def validate_data(data: Dict[str, Any], config: Dict[str, Any], template_type: Optional[str] = None) -> List[str]:
    """
    Validate that required fields are present in data.

    Args:
        data: Data dictionary
        config: Template configuration dictionary
        template_type: Template type (e.g., "proposal")

    Returns:
        List of warning messages for missing fields
    """
    warnings_list = []

    if not config or not template_type:
        return warnings_list

    # Get template type config
    templates = config.get('templates', {})
    template_config = templates.get(template_type, {})

    if not template_config:
        return warnings_list

    # Check required fields
    required_fields = template_config.get('required_fields', [])
    for field in required_fields:
        if field not in data:
            warnings_list.append(f"Missing required field: {field}")

    return warnings_list


def resolve_value(data: Dict[str, Any], key: str) -> Any:
    """
    Resolve a placeholder key to its value, supporting nested keys.

    Args:
        data: Data dictionary
        key: Placeholder key (e.g., "company_name" or "items.0.name")

    Returns:
        Resolved value or None if not found
    """
    keys = key.split('.')
    value = data

    for k in keys:
        if isinstance(value, dict):
            value = value.get(k)
        elif isinstance(value, list):
            try:
                index = int(k)
                if 0 <= index < len(value):
                    value = value[index]
                else:
                    return None
            except ValueError:
                return None
        else:
            return None

        if value is None:
            return None

    return value


def format_value(value: Any) -> str:
    """
    Format a value for insertion into template.

    Args:
        value: Value to format

    Returns:
        Formatted string
    """
    if value is None:
        return ""
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, bool):
        return "Yes" if value else "No"
    return str(value)


def replace_placeholder_docx(
    doc: Any,
    placeholder: str,
    value: str,
    preserve_formatting: bool = True
) -> int:
    """
    Replace a placeholder in a Word document, preserving formatting.

    Args:
        doc: Document object
        placeholder: Placeholder key (e.g., "company_name")
        value: Replacement value
        preserve_formatting: Whether to preserve formatting

    Returns:
        Number of replacements made
    """
    if Document is None:
        raise ImportError("python-docx is required for .docx files. Install with: pip install python-docx")

    marker = f"{{{{{placeholder}}}}}"
    count = 0

    # Replace in paragraphs
    for paragraph in doc.paragraphs:
        if marker in paragraph.text:
            count += _replace_in_paragraph(paragraph, marker, value, preserve_formatting)

    # Replace in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if marker in paragraph.text:
                        count += _replace_in_paragraph(paragraph, marker, value, preserve_formatting)

    # Replace in headers and footers
    for section in doc.sections:
        # Header
        header = section.header
        for paragraph in header.paragraphs:
            if marker in paragraph.text:
                count += _replace_in_paragraph(paragraph, marker, value, preserve_formatting)

        # Footer
        footer = section.footer
        for paragraph in footer.paragraphs:
            if marker in paragraph.text:
                count += _replace_in_paragraph(paragraph, marker, value, preserve_formatting)

    return count


def _replace_in_paragraph(paragraph, marker: str, value: str, preserve_formatting: bool) -> int:
    """
    Replace marker in a paragraph, preserving formatting of the placeholder text.

    Args:
        paragraph: Paragraph object
        marker: Placeholder marker (e.g., "{{company_name}}")
        value: Replacement value
        preserve_formatting: Whether to preserve formatting

    Returns:
        Number of replacements made (0 or 1)
    """
    if marker not in paragraph.text:
        return 0

    # Find the run containing the marker
    full_text = paragraph.text
    marker_start = full_text.find(marker)

    if marker_start == -1:
        return 0

    # Calculate which run contains the marker
    current_pos = 0
    target_run = None

    for run in paragraph.runs:
        run_end = current_pos + len(run.text)
        if current_pos <= marker_start < run_end:
            target_run = run
            break
        current_pos = run_end

    if target_run is None:
        return 0

    # Get formatting from the run if preserving
    if preserve_formatting:
        font_name = target_run.font.name
        font_size = target_run.font.size
        font_bold = target_run.font.bold
        font_italic = target_run.font.italic
        font_underline = target_run.font.underline
        font_color = target_run.font.color
    else:
        font_name = None
        font_size = None
        font_bold = None
        font_italic = None
        font_underline = None
        font_color = None

    # Replace the marker in the run text
    run_text = target_run.text
    new_text = run_text.replace(marker, value)
    target_run.text = new_text

    # Restore formatting if preserving
    if preserve_formatting:
        if font_name:
            target_run.font.name = font_name
        if font_size:
            target_run.font.size = font_size
        if font_bold is not None:
            target_run.font.bold = font_bold
        if font_italic is not None:
            target_run.font.italic = font_italic
        if font_underline is not None:
            target_run.font.underline = font_underline
        if font_color:
            target_run.font.color = font_color

    return 1


def replicate_table_rows_docx(doc: Any, list_key: str, list_data: List[Dict[str, Any]]) -> None:
    """
    Replicate table rows for list data in Word document.

    Looks for a row with {{list_key.0.field}} placeholders and replicates it for each item.

    Args:
        doc: Document object
        list_key: Key of the list in data (e.g., "items")
        list_data: List of dictionaries to replicate
    """
    if not list_data:
        return

    # Pattern to match list placeholders like {{items.0.name}}
    escaped_key = re.escape(list_key)
    pattern_str = r'\{\{' + escaped_key + r'\.0\.([^}]+)\}\}'
    pattern = re.compile(pattern_str)

    for table in doc.tables:
        # Find the template row (contains {{list_key.0.field}})
        template_row_idx = None
        for idx, row in enumerate(table.rows):
            row_text = ' '.join(cell.text for cell in row.cells)
            if pattern.search(row_text):
                template_row_idx = idx
                break

        if template_row_idx is None:
            continue

        # Get template row
        template_row = table.rows[template_row_idx]
        template_cells = [cell for cell in template_row.cells]

        # Insert new rows after template row
        for i, item in enumerate(list_data):
            if i == 0:
                # Use template row for first item
                row = template_row
            else:
                # Insert new row
                row = table.add_row()
                # Copy cell structure from template
                for j, template_cell in enumerate(template_cells):
                    if j < len(row.cells):
                        # Copy paragraphs from template cell
                        for para in template_cell.paragraphs:
                            new_para = row.cells[j].add_paragraph()
                            for run in para.runs:
                                new_run = new_para.add_run(run.text)
                                # Copy formatting
                                new_run.font.name = run.font.name
                                if run.font.size:
                                    new_run.font.size = run.font.size
                                new_run.font.bold = run.font.bold
                                new_run.font.italic = run.font.italic

            # Replace placeholders in this row
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            # Replace {{list_key.0.field}} with {{list_key.i.field}}
                            new_text = pattern.sub(
                                lambda m: f"{{{{{list_key}.{i}.{m.group(1)}}}}}",
                                run.text
                            )
                            run.text = new_text

        # Replace all placeholders in the replicated rows
        for i, item in enumerate(list_data):
            row_idx = template_row_idx + i
            if row_idx < len(table.rows):
                row = table.rows[row_idx]
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                # Replace {{list_key.i.field}} with actual values
                                for key, val in item.items():
                                    marker = f"{{{{{list_key}.{i}.{key}}}}}"
                                    if marker in run.text:
                                        run.text = run.text.replace(marker, format_value(val))


def fill_template_docx(template_path: Path, data: Dict[str, Any], output_path: Path) -> None:
    """
    Fill a Word document template with data.

    Args:
        template_path: Path to template .docx file
        data: Data dictionary
        output_path: Path to output .docx file
    """
    if Document is None:
        raise ImportError("python-docx is required for .docx files. Install with: pip install python-docx")

    doc = Document(str(template_path))

    # First pass: handle list replication
    for key, value in data.items():
        if isinstance(value, list) and value and isinstance(value[0], dict):
            replicate_table_rows_docx(doc, key, value)

    # Second pass: replace simple placeholders
    for key, value in data.items():
        if not isinstance(value, list):
            resolved_value = resolve_value(data, key)
            if resolved_value is not None:
                formatted_value = format_value(resolved_value)
                replace_placeholder_docx(doc, key, formatted_value)

    # Third pass: replace nested placeholders (e.g., items.0.name)
    all_placeholders = set()
    for paragraph in doc.paragraphs:
        for match in PLACEHOLDER_PATTERN.finditer(paragraph.text):
            placeholder = match.group(1)
            # Skip IMAGE placeholders
            if not placeholder.startswith('IMAGE:'):
                all_placeholders.add(placeholder)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for match in PLACEHOLDER_PATTERN.finditer(paragraph.text):
                        placeholder = match.group(1)
                        if not placeholder.startswith('IMAGE:'):
                            all_placeholders.add(placeholder)

    for placeholder in all_placeholders:
        resolved_value = resolve_value(data, placeholder)
        if resolved_value is not None:
            formatted_value = format_value(resolved_value)
            replace_placeholder_docx(doc, placeholder, formatted_value)

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    print(f"✓ Template filled: {output_path.name}")


def replicate_table_rows_xlsx(workbook: Any, list_key: str, list_data: List[Dict[str, Any]], config: Dict[str, Any]) -> None:
    """
    Replicate table rows for list data in Excel workbook.

    Args:
        workbook: OpenPyXL workbook object
        list_key: Key of the list in data (e.g., "items")
        list_data: List of dictionaries to replicate
        config: Template configuration (for formula preservation)
    """
    if not list_data:
        return

    # Pattern to match list placeholders like {{items.0.name}}
    escaped_key = re.escape(list_key)
    pattern_str = r'\{\{' + escaped_key + r'\.0\.([^}]+)\}\}'
    pattern = re.compile(pattern_str)

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]

        # Find template row (contains {{list_key.0.field}})
        template_row_idx = None
        for idx, row in enumerate(sheet.iter_rows(), start=1):
            row_text = ' '.join(str(cell.value) if cell.value else '' for cell in row)
            if pattern.search(row_text):
                template_row_idx = idx
                break

        if template_row_idx is None:
            continue

        # Get template row
        template_row = list(sheet.iter_rows(min_row=template_row_idx, max_row=template_row_idx))[0]

        # Check which cells are formulas (from config)
        formula_cells = set()
        if config:
            templates = config.get('templates', {})
            for template_config in templates.values():
                sheets_config = template_config.get('sheets', {})
                sheet_config = sheets_config.get(sheet_name, {})
                formulas = sheet_config.get('formulas', [])
                for formula_cell in formulas:
                    if isinstance(formula_cell, str):
                        # Convert cell reference to column, row
                        col_letter = ''.join(c for c in formula_cell if c.isalpha())
                        col_idx = openpyxl.utils.column_index_from_string(col_letter)
                        formula_cells.add((col_idx, template_row_idx))

        # Insert new rows after template row
        for i, item in enumerate(list_data):
            if i == 0:
                # Use template row for first item
                row_idx = template_row_idx
            else:
                # Insert new row
                sheet.insert_rows(template_row_idx + i)
                row_idx = template_row_idx + i

            # Copy values and formatting from template row
            for col_idx, template_cell in enumerate(template_row, start=1):
                new_cell = sheet.cell(row=row_idx, column=col_idx)

                # Check if this is a formula cell (preserve formula)
                if (col_idx, template_row_idx) in formula_cells:
                    # Copy formula, adjusting row references
                    if template_cell.data_type == 'f':
                        formula = template_cell.value
                        if formula and isinstance(formula, str) and formula.startswith('='):
                            # Adjust row references in formula
                            adjusted_formula = re.sub(
                                r'(\$?)([A-Z]+)(\$?)(\d+)',
                                lambda m: f"{m.group(1)}{m.group(2)}{m.group(3)}{int(m.group(4)) + i}",
                                formula
                            )
                            new_cell.value = adjusted_formula
                        else:
                            new_cell.value = template_cell.value
                    else:
                        new_cell.value = template_cell.value
                else:
                    # Regular cell: replace placeholder
                    cell_value = str(template_cell.value) if template_cell.value else ""
                    # Replace {{list_key.0.field}} with {{list_key.i.field}}
                    new_value = pattern.sub(
                        lambda m: f"{{{{{list_key}.{i}.{m.group(1)}}}}}",
                        cell_value
                    )
                    new_cell.value = new_value

                # Copy formatting
                if template_cell.has_style:
                    new_cell.font = template_cell.font
                    new_cell.fill = template_cell.fill
                    new_cell.border = template_cell.border
                    new_cell.alignment = template_cell.alignment
                    new_cell.number_format = template_cell.number_format

        # Replace placeholders in replicated rows
        for i, item in enumerate(list_data):
            row_idx = template_row_idx + i
            for col_idx, cell in enumerate(list(sheet.iter_rows(min_row=row_idx, max_row=row_idx))[0], start=1):
                if cell.value and isinstance(cell.value, str):
                    # Replace {{list_key.i.field}} with actual values
                    for key, val in item.items():
                        marker = f"{{{{{list_key}.{i}.{key}}}}}"
                        if marker in str(cell.value):
                            cell.value = str(cell.value).replace(marker, format_value(val))


def fill_template_xlsx(template_path: Path, data: Dict[str, Any], output_path: Path, config: Optional[Dict[str, Any]] = None) -> None:
    """
    Fill an Excel template with data.

    Args:
        template_path: Path to template .xlsx file
        data: Data dictionary
        output_path: Path to output .xlsx file
        config: Template configuration (optional)
    """
    if openpyxl is None:
        raise ImportError("openpyxl is required for .xlsx files. Install with: pip install openpyxl")

    workbook = openpyxl.load_workbook(str(template_path))

    # First pass: handle list replication
    for key, value in data.items():
        if isinstance(value, list) and value and isinstance(value[0], dict):
            replicate_table_rows_xlsx(workbook, key, value, config or {})

    # Second pass: replace simple placeholders
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]

        # Get formula cells from config (don't replace these)
        formula_cells = set()
        if config:
            templates = config.get('templates', {})
            for template_config in templates.values():
                sheets_config = template_config.get('sheets', {})
                sheet_config = sheets_config.get(sheet_name, {})
                formulas = sheet_config.get('formulas', [])
                for formula_cell in formulas:
                    if isinstance(formula_cell, str):
                        col_letter = ''.join(c for c in formula_cell if c.isalpha())
                        col_idx = openpyxl.utils.column_index_from_string(col_letter)
                        row_num = int(''.join(c for c in formula_cell if c.isdigit()))
                        formula_cells.add((col_idx, row_num))

        for row in sheet.iter_rows():
            for cell in row:
                # Skip formula cells
                if (cell.column, cell.row) in formula_cells:
                    continue

                if cell.value and isinstance(cell.value, str):
                    cell_text = str(cell.value)

                    # Skip IMAGE placeholders
                    if IMAGE_PLACEHOLDER_PATTERN.search(cell_text):
                        continue

                    # Replace placeholders
                    for match in PLACEHOLDER_PATTERN.finditer(cell_text):
                        placeholder = match.group(1)
                        if placeholder.startswith('IMAGE:'):
                            continue

                        resolved_value = resolve_value(data, placeholder)
                        if resolved_value is not None:
                            formatted_value = format_value(resolved_value)
                            cell_text = cell_text.replace(match.group(0), formatted_value)

                    cell.value = cell_text

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(str(output_path))
    print(f"✓ Template filled: {output_path.name}")


def fill_template_pptx(template_path: Path, data: Dict[str, Any], output_path: Path) -> None:
    """
    Fill a PowerPoint template with data.

    Args:
        template_path: Path to template .pptx file
        data: Data dictionary
        output_path: Path to output .pptx file
    """
    if Presentation is None:
        raise ImportError("python-pptx is required for .pptx files. Install with: pip install python-pptx")

    prs = Presentation(str(template_path))

    # Replace placeholders in all slides
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                # Replace in paragraphs
                for paragraph in text_frame.paragraphs:
                    if paragraph.text:
                        # Skip IMAGE placeholders
                        if IMAGE_PLACEHOLDER_PATTERN.search(paragraph.text):
                            continue

                        # Replace placeholders
                        for match in PLACEHOLDER_PATTERN.finditer(paragraph.text):
                            placeholder = match.group(1)
                            if placeholder.startswith('IMAGE:'):
                                continue

                            resolved_value = resolve_value(data, placeholder)
                            if resolved_value is not None:
                                formatted_value = format_value(resolved_value)

                                # Preserve formatting: find the run containing the placeholder
                                full_text = paragraph.text
                                marker_start = full_text.find(match.group(0))

                                if marker_start != -1:
                                    current_pos = 0
                                    for run in paragraph.runs:
                                        run_end = current_pos + len(run.text)
                                        if current_pos <= marker_start < run_end:
                                            # Replace in this run
                                            run.text = run.text.replace(match.group(0), formatted_value)
                                            break
                                        current_pos = run_end

            # Handle table shapes
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            if paragraph.text:
                                # Skip IMAGE placeholders
                                if IMAGE_PLACEHOLDER_PATTERN.search(paragraph.text):
                                    continue

                                # Replace placeholders
                                for match in PLACEHOLDER_PATTERN.finditer(paragraph.text):
                                    placeholder = match.group(1)
                                    if placeholder.startswith('IMAGE:'):
                                        continue

                                    resolved_value = resolve_value(data, placeholder)
                                    if resolved_value is not None:
                                        formatted_value = format_value(resolved_value)

                                        # Replace in runs
                                        for run in paragraph.runs:
                                            if match.group(0) in run.text:
                                                run.text = run.text.replace(match.group(0), formatted_value)

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✓ Template filled: {output_path.name}")


def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Fill Word, Excel, or PowerPoint templates with data from JSON",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Direct mode: specify template and data files
  %(prog)s --template templates/word/proposal_template.docx --data data.json --output output/proposal.docx

  # Config mode: use template_config.yaml
  %(prog)s --config templates/template_config.yaml --type proposal --data data.json

  # Auto-detect format from template extension
  %(prog)s --template templates/excel/pricing.xlsx --data data.json --output output/pricing.xlsx
        """
    )

    # Mode selection
    mode_group = parser.add_mutually_exclusive_group(required=True)
    mode_group.add_argument(
        "--template",
        type=Path,
        help="Path to template file (.docx, .xlsx, or .pptx)"
    )
    mode_group.add_argument(
        "--config",
        type=Path,
        help="Path to template_config.yaml file"
    )

    parser.add_argument(
        "--type",
        type=str,
        help="Template type (required when using --config)"
    )
    parser.add_argument(
        "--data",
        type=Path,
        required=True,
        help="Path to JSON data file"
    )
    parser.add_argument(
        "--output",
        type=Path,
        help="Path to output file (required in direct mode, optional in config mode)"
    )

    args = parser.parse_args()

    # Load data
    try:
        data = load_data(args.data)
    except Exception as e:
        print(f"Error loading data file: {e}", file=sys.stderr)
        sys.exit(1)

    # Determine template path and output path
    if args.template:
        # Direct mode
        template_path = args.template
        if not args.output:
            print("Error: --output is required when using --template", file=sys.stderr)
            sys.exit(1)
        output_path = args.output
        config = None
    else:
        # Config mode
        if not args.type:
            print("Error: --type is required when using --config", file=sys.stderr)
            sys.exit(1)

        try:
            config = load_template_config(args.config)
        except Exception as e:
            print(f"Error loading config file: {e}", file=sys.stderr)
            sys.exit(1)

        # Get template path from config
        templates = config.get('templates', {})
        template_config = templates.get(args.type, {})
        if not template_config:
            print(f"Error: Template type '{args.type}' not found in config", file=sys.stderr)
            sys.exit(1)

        template_path = Path(template_config.get('template_path', ''))
        if not template_path.exists():
            # Try relative to config file
            template_path = args.config.parent / template_path
            if not template_path.exists():
                print(f"Error: Template file not found: {template_path}", file=sys.stderr)
                sys.exit(1)

        # Get output path from config or use default
        if args.output:
            output_path = args.output
        else:
            output_name = template_config.get('output_name', template_path.stem + '_filled' + template_path.suffix)
            output_path = Path('output') / output_name

    # Validate template exists
    if not template_path.exists():
        print(f"Error: Template file not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    # Validate data (warn but continue)
    if config:
        warnings_list = validate_data(data, config, args.type)
        for warning in warnings_list:
            warnings.warn(warning, UserWarning)

    # Auto-detect format from extension
    ext = template_path.suffix.lower()

    try:
        if ext == ".docx":
            fill_template_docx(template_path, data, output_path)
        elif ext == ".xlsx":
            fill_template_xlsx(template_path, data, output_path, config)
        elif ext == ".pptx":
            fill_template_pptx(template_path, data, output_path)
        else:
            print(f"Error: Unsupported file format: {ext}. Supported formats: .docx, .xlsx, .pptx", file=sys.stderr)
            sys.exit(1)

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
